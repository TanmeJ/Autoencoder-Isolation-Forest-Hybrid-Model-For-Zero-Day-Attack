from pathlib import Path
import re
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path(r"C:\Users\jtanm\OneDrive\Documents\cn\Network_Anomaly_Detection_Presentation.pptx")
OUTPUT = SOURCE.with_name("final_presentation.pptx")

BG = RGBColor(0xF8, 0xFA, 0xFC)
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
ACCENT = RGBColor(0x10, 0xB9, 0x81)
TEXT = RGBColor(0x0F, 0x17, 0x2A)
SECONDARY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Poppins"
BODY_FONT = "Inter"
MAX_BULLETS = 4
MAX_WORDS = 7


def apply_fill(target, color):
    target.solid()
    target.fore_color.rgb = color


def set_shape_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, width=1.25):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, text, x, y, w, h, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = TITLE_FONT if bold else BODY_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                parts.append(text)
    return parts


def clean_line(line):
    line = re.sub(r"^[\-\u2022\*\d\.\)\s]+", "", line.strip())
    line = re.sub(r"\s+", " ", line)
    return line.strip(" .")


def compact_line(text):
    words = clean_line(text).split()
    if len(words) <= MAX_WORDS:
        return " ".join(words)
    return " ".join(words[:MAX_WORDS])


def paragraph_to_bullets(text):
    candidates = []
    for raw in re.split(r"[\n\r]+|(?<=[.!?])\s+", text):
        line = clean_line(raw)
        if line:
            candidates.append(compact_line(line))
    return candidates[:MAX_BULLETS]


def is_path_text(text):
    return bool(re.match(r"^[A-Za-z]:\\", text.strip()))


def add_bullets(slide, bullets, x, y, w, h):
    for index, bullet in enumerate(bullets[:MAX_BULLETS]):
        row_y = y + Inches(index * 0.72)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, row_y + Inches(0.16), Inches(0.12), Inches(0.12))
        set_shape_fill(dot, ACCENT if index % 2 else PRIMARY)
        set_no_line(dot)
        add_textbox(slide, bullet, x + Inches(0.28), row_y, w - Inches(0.28), Inches(0.38), 19, TEXT)


def add_header(slide, title):
    add_textbox(slide, title, Inches(0.75), Inches(0.58), Inches(5.7), Inches(0.55), 25, TEXT, True)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.22), Inches(0.82), Inches(0.06))
    set_shape_fill(bar, PRIMARY)
    set_no_line(bar)


def add_icon_circle(slide, x, y, label, fill=PRIMARY):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.58), Inches(0.58))
    set_shape_fill(outer, fill)
    set_no_line(outer)
    add_textbox(slide, label, x, y + Inches(0.08), Inches(0.58), Inches(0.34), 15, WHITE, True, PP_ALIGN.CENTER)


def add_visual_panel(slide, index):
    x, y, w, h = Inches(7.0), Inches(1.45), Inches(5.4), Inches(4.7)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_shape_fill(panel, WHITE)
    set_line(panel, RGBColor(0xE2, 0xE8, 0xF0), 1)

    if index == 5:
        labels = [("TN", "35.9K", PRIMARY), ("FP", "1.9K", ACCENT), ("FN", "388", ACCENT), ("TP", "5", PRIMARY)]
        for i, (name, value, color) in enumerate(labels):
            col = i % 2
            row = i // 2
            bx = x + Inches(0.48 + col * 2.45)
            by = y + Inches(0.65 + row * 1.75)
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, bx, by, Inches(1.95), Inches(1.24))
            set_shape_fill(cell, RGBColor(0xF8, 0xFA, 0xFC))
            set_line(cell, RGBColor(0xE2, 0xE8, 0xF0), 1)
            add_textbox(slide, name, bx + Inches(0.18), by + Inches(0.18), Inches(1.5), Inches(0.26), 12, SECONDARY, True)
            add_textbox(slide, value, bx + Inches(0.18), by + Inches(0.5), Inches(1.5), Inches(0.38), 22, color, True)
        return

    for i, height in enumerate([2.45, 1.85, 2.95, 2.2]):
        bx = x + Inches(0.65 + i * 1.0)
        by = y + Inches(3.55 - height * 0.45)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bx, by, Inches(0.42), Inches(height * 0.45))
        set_shape_fill(bar, PRIMARY if i % 2 == 0 else ACCENT)
        set_no_line(bar)

    for i, label in enumerate(["Data", "Train", "Score"]):
        add_icon_circle(slide, x + Inches(0.75 + i * 1.55), y + Inches(0.7), str(i + 1), ACCENT if i == 2 else PRIMARY)
        add_textbox(slide, label, x + Inches(0.52 + i * 1.55), y + Inches(1.35), Inches(1.05), Inches(0.26), 12, SECONDARY, True, PP_ALIGN.CENTER)


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def set_background(slide):
    apply_fill(slide.background.fill, BG)


def build_title_slide(slide, title, subtitle, footer):
    set_background(slide)
    add_textbox(slide, title, Inches(1.35), Inches(2.15), Inches(10.65), Inches(0.8), 38, TEXT, True, PP_ALIGN.CENTER)
    add_textbox(slide, subtitle, Inches(2.2), Inches(3.08), Inches(8.95), Inches(0.48), 19, SECONDARY, False, PP_ALIGN.CENTER)
    add_textbox(slide, footer, Inches(2.6), Inches(5.7), Inches(8.1), Inches(0.3), 11, SECONDARY, False, PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(5.48), Inches(3.82), Inches(2.35), Inches(0.07))
    set_shape_fill(line, ACCENT)
    set_no_line(line)


def build_section_slide(slide, title, bullets):
    set_background(slide)
    add_textbox(slide, title, Inches(1.0), Inches(2.45), Inches(11.33), Inches(0.85), 36, TEXT, True, PP_ALIGN.CENTER)
    if bullets:
        add_textbox(slide, bullets[0], Inches(2.25), Inches(3.55), Inches(8.8), Inches(0.36), 17, SECONDARY, False, PP_ALIGN.CENTER)


def build_content_slide(slide, title, bullets, index):
    set_background(slide)
    add_header(slide, title)
    add_bullets(slide, bullets, Inches(0.95), Inches(1.68), Inches(5.55), Inches(4.6))
    add_visual_panel(slide, index)


def add_fade_transitions(pptx_path):
    slide_pattern = re.compile(r"ppt/slides/slide\d+\.xml$")
    temp_path = pptx_path.with_suffix(".tmp.pptx")
    transition_xml = '<p:transition spd="med"><p:fade/></p:transition>'

    with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if slide_pattern.match(item.filename):
                xml = data.decode("utf-8")
                xml = re.sub(r"<p:transition\b.*?</p:transition>", "", xml, flags=re.DOTALL)
                if "</p:clrMapOvr>" in xml:
                    xml = xml.replace("</p:clrMapOvr>", f"</p:clrMapOvr>{transition_xml}", 1)
                else:
                    xml = xml.replace("</p:cSld>", f"</p:cSld>{transition_xml}", 1)
                data = xml.encode("utf-8")
            target.writestr(item, data)
    temp_path.replace(pptx_path)


def modernize():
    prs = Presentation(SOURCE)

    for idx, slide in enumerate(prs.slides):
        parts = [part for part in slide_text(slide) if not is_path_text(part)]
        title = clean_line(parts[0]) if parts else f"Slide {idx + 1}"
        body = "\n".join(parts[1:]) if len(parts) > 1 else ""
        bullets = paragraph_to_bullets(body)
        clear_slide(slide)

        if idx == 0:
            subtitle = bullets[0] if bullets else "Hybrid ML system"
            build_title_slide(slide, title, subtitle, str(SOURCE.parent))
        elif title.lower() in {"introduction", "conclusion"}:
            build_section_slide(slide, title, bullets)
        else:
            build_content_slide(slide, title, bullets, idx)

    prs.save(OUTPUT)
    add_fade_transitions(OUTPUT)
    print(f"Saved {OUTPUT}")


if __name__ == "__main__":
    modernize()
